from pptx import Presentation
from datetime import date
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

weekly_ppt = Presentation('Weekly Price Index.pptx')
slide = weekly_ppt.slides[7]
shapes = slide.shapes


curr = open('curr.txt', 'r')
current = curr.readline()
past = open('past.txt', 'r')
pastline = past.readline()
pastline = pastline[:-2]
today = date.today()

for shape in shapes:
    if shape.shape_type == 1:
        if shape.text[0:5] == "Steel":
            text = shape

text.text="Steel (U.S) \n2W- " + pastline
pastline = past.readline()
pastline = pastline[:-2]
text.text = text.text + " ($" + pastline + "/T = $" + str(round(float(pastline)*0.12,2)) + "/w)" + "to\n" + today.strftime("%B %d") + "($" + current + "/T = $" + str(round(float(current)*0.12,2)) + "/w) \n"
if (float(current)-float(pastline))/float(pastline)<0:
    text.text=text.text+str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n\n"
if (float(current)-float(pastline))/float(pastline)>0:
    text.text=text.text+str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n\n"
if (float(current)-float(pastline))/float(pastline)==0:
    text.text + "no change\n\n"
pastline=past.readline()
pastline = pastline[:-2]

p = text.text_frame.add_paragraph()
run = p.add_run()
run.text = "Last meeting: "
run.font.name = 'Calibri'
run.font.size = Pt(18)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4

run = p.add_run()
run.text = pastline+"\n\n"
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2

pastline=past.readline()
pastline = pastline[:-2]
text.text = text.text+"YTD- " + pastline
pastline = past.readline()
pastline = pastline[:-2]
text.text = text.text + " ($" + pastline + "/T = $" + str(round(float(pastline)*0.12,2)) + "/w)" + "to\n" + today.strftime("%B %d") + "($" + current + "/T = $" + str(round(float(current)*0.12,2)) + "/w) \n"
if (float(current)-float(pastline))/float(pastline)<0:
    text.text=text.text+str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n\n"
if (float(current)-float(pastline))/float(pastline)>0:
    text.text=text.text+str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n\n"
if (float(current)-float(pastline))/float(pastline)==0:
    text.text + "no change\n\n"
pastline=past.readline()
pastline = pastline[:-2]
text.text = text.text+"1Y- " + pastline
pastline = past.readline()
pastline = pastline[:-2]
text.text = text.text + " ($" + pastline + "/T = $" + str(round(float(pastline)*0.12,2)) + "/w)" + "to\n" + today.strftime("%B %d") + "($" + current + "/T = $" + str(round(float(current)*0.12,2)) + "/w) \n"
if (float(current)-float(pastline))/float(pastline)<0:
    text.text=text.text+str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n\n"
if (float(current)-float(pastline))/float(pastline)>0:
    text.text=text.text+str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n\n"
if (float(current)-float(pastline))/float(pastline)==0:
    text.text + "no change\n\n"
