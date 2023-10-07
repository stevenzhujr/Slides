from pptx import Presentation
from pptx.util import Inches
# slide 7
weekly_ppt = Presentation('Weekly Price Index.pptx')
slide = weekly_ppt.slides[6]
shapes = slide.shapes

for shape in shapes:
    if shape.shape_type == 13:
        shapes.element.remove(shape.element)

slide = weekly_ppt.slides[7]
shapes = slide.shapes
# slide 8
for shape in shapes:
    if shape.shape_type == 13:
        shapes.element.remove(shape.element)

weekly_ppt.save('Weekly Price Index.pptx')