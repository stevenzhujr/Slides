from pptx import Presentation
from pptx.util import Inches
from datetime import date
from pptx.util import Pt
from pptx.enum.dml import MSO_THEME_COLOR

# Setup

curr = open('curr.txt', 'r')
current = curr.readline()
past = open('past.txt', 'r')
pastline = past.readline()
pastline = pastline[:-2]
today = date.today()

# Update slides

weekly_ppt = Presentation('Weekly Price Index.pptx')
slide = weekly_ppt.slides[0]
shapes = slide.shapes
#region
# slide 7
slide = weekly_ppt.slides[6]
ref_element = slide.shapes[0]._element
shapes = slide.shapes

pic='SevenHeader.png'
pica=slide.shapes.add_picture(pic, Inches(2.0), Inches(1.0), width=Inches(7.0))
ref_element.addprevious(pica._element)

pic='SevenTop.png'
pica=slide.shapes.add_picture(pic, Inches(2.0), Inches(1.8), width=Inches(7.0))
ref_element.addprevious(pica._element)

pic='SevenMiddle.png'
pica=slide.shapes.add_picture(pic, Inches(2.0), Inches(3.0), width=Inches(7.0))
ref_element.addprevious(pica._element)

pic='SevenBottom.png'
pica=slide.shapes.add_picture(pic, Inches(2.0), Inches(5.5), width=Inches(7.0))
ref_element.addprevious(pica._element)
#endregion
#region
# slide 8
slide = weekly_ppt.slides[7]
shapes = slide.shapes

pic='SevenHeader.png'
pica=slide.shapes.add_picture(pic, Inches(0.7), Inches(1.0), width=Inches(6.0))

pic='EightTop.png'
pica=slide.shapes.add_picture(pic, Inches(0.7), Inches(1.8), width=Inches(6.0))

pic='EightMiddle.png'
pica=slide.shapes.add_picture(pic, Inches(0.7), Inches(2.5), width=Inches(6.0))

pic='EightBottom.png'
pica=slide.shapes.add_picture(pic, Inches(0.7), Inches(4.0), width=Inches(6.0))

pic='EightRight.png'
pica=slide.shapes.add_picture(pic, Inches(7.0), Inches(5.3), width=Inches(6.0))

for shape in shapes:
    if shape.shape_type == 1:
        if shape.text[0:5]=="Steel":
            text = shape

text.text_frame.clear()

text.text_frame.paragraphs[0].text = "Steel (U.S)"
text.text_frame.paragraphs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
text.text_frame.paragraphs[0].font.bold=True
p = text.text_frame.add_paragraph()
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.bold=True
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = "2W- "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = pastline
pastline = past.readline()
pastline = pastline[:-2]
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = " ($" + pastline + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(pastline)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ") to\n" + today.strftime("%B %d") + " ($" + current + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(current)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ")"
p = text.text_frame.add_paragraph()
p.font.name = 'Calibri'
p.font.size = Pt(16)
p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
p.font.bold=True
if (float(current)-float(pastline))/float(pastline)<0:
    p.text=str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n"
if (float(current)-float(pastline))/float(pastline)>0:
    p.text=str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n"
if (float(current)-float(pastline))/float(pastline)==0:
    p.text="no change\n"
pastline=past.readline()
pastline = pastline[:-2]

p = text.text_frame.add_paragraph()
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = 'Last meeting: '
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.bold=True
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
run.text = pastline+"\n"
pastline=past.readline()
pastline = pastline[:-2]

p = text.text_frame.add_paragraph()
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.bold=True
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = "YTD- "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = pastline
pastline = past.readline()
pastline = pastline[:-2]
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = " ($" + pastline + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(pastline)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ") to\n" + today.strftime("%B %d") + " ($" + current + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(current)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ")"
p = text.text_frame.add_paragraph()
p.font.name = 'Calibri'
p.font.size = Pt(16)
p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
p.font.bold=True
if (float(current)-float(pastline))/float(pastline)<0:
    p.text=str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n"
if (float(current)-float(pastline))/float(pastline)>0:
    p.text=str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n"
if (float(current)-float(pastline))/float(pastline)==0:
    p.text="no change\n"
pastline=past.readline()
pastline = pastline[:-2]

p = text.text_frame.add_paragraph()
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.bold=True
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = "1Y- "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = pastline
pastline = past.readline()
pastline = pastline[:-2]
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = " ($" + pastline + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(pastline)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ") to\n" + today.strftime("%B %d") + " ($" + current + "/T = "
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.font.bold=True
run.text = "$" + str(round(float(current)*0.12,2)) + "/w"
run = p.add_run()
run.font.name = 'Calibri'
run.font.size = Pt(16)
run.font.color.theme_color = MSO_THEME_COLOR.ACCENT_4
run.text = ")"
p = text.text_frame.add_paragraph()
p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_2
p.font.bold=True
p.font.name = 'Calibri'
p.font.size = Pt(16)
if (float(current)-float(pastline))/float(pastline)<0:
    p.text=str(round(-1*(float(current)-float(pastline))/float(pastline)*100,2))+"% decrease\n"
if (float(current)-float(pastline))/float(pastline)>0:
    p.text=str(round((float(current)-float(pastline))/float(pastline)*100,2))+"% increase\n"
if (float(current)-float(pastline))/float(pastline)==0:
    p.text="no change\n"
#endregion

weekly_ppt.save('Weekly Price Index.pptx')